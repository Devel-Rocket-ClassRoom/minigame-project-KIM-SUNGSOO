# -*- coding: utf-8 -*-
"""
SwordRush TD 발표자료 v2 — 기존 PPT 와 같은 구조 + Build 3+ 폴리시 슬라이드 추가.

색상 팔레트 (원본 동일):
  Navy   #1E2761  primary
  Gold   #F59E0B  accent
  Card   #FFFFFF  white card bg
  LBlue  #EEF3FB  light blue panel
  Code   #0F172A  code block bg
  Red    #DC2626  warning / sell
  Ice    #99C5F0  frost tint accent (신규)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY   = RGBColor(0x1E, 0x27, 0x61)
GOLD   = RGBColor(0xF5, 0x9E, 0x0B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE  = RGBColor(0xEE, 0xF3, 0xFB)
CODE   = RGBColor(0x0F, 0x17, 0x2A)
RED    = RGBColor(0xDC, 0x26, 0x26)
ICE    = RGBColor(0x99, 0xC5, 0xF0)
GRAY7  = RGBColor(0x33, 0x33, 0x33)
GRAY5  = RGBColor(0x55, 0x55, 0x55)
GRAY9  = RGBColor(0x99, 0x99, 0x99)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)

# ---------------------------------------------------------------
# 헬퍼 함수
# ---------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, size=14, bold=False, color=GRAY7, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font='맑은 고딕', italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

def add_footer(slide, page, total=15):
    add_text(slide, 0.50, 5.26, 6.00, 0.25, 'Unity 2D Tower Defense  ·  김성수',
             size=9, color=GRAY9)
    add_text(slide, 8.80, 5.26, 0.70, 0.25, f'{page} / {total}',
             size=9, color=GRAY9, align=PP_ALIGN.RIGHT)

def add_section_header(slide, eyebrow, title):
    add_text(slide, 0.50, 0.32, 9.00, 0.30, eyebrow, size=11, bold=True, color=GOLD)
    add_text(slide, 0.50, 0.58, 9.00, 0.70, title, size=24, bold=True, color=NAVY)

# ---------------------------------------------------------------
# Presentation init — 16:9
# ---------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank = prs.slide_layouts[6]
TOTAL = 19

# ---------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 10, 5.625, WHITE)
add_rect(s, 0, 0, 0.18, 5.625, GOLD)
add_text(s, 0.80, 1.55, 8.00, 0.40, 'UNITY 2D MOBILE GAME PROJECT',
         size=12, bold=True, color=GOLD)
add_text(s, 0.80, 1.95, 9.00, 1.00, 'Sword Rush TD',
         size=54, bold=True, color=NAVY)
add_text(s, 0.80, 3.05, 9.00, 0.50,
         '웨이브를 막아라 — 타워 · 영웅 · 특수능력으로 짜는 디펜스 전략',
         size=15, color=GRAY5)
add_rect(s, 0.80, 3.85, 1.60, 0.05, GOLD)
add_text(s, 0.80, 4.05, 9.00, 0.35, '발표자  김성수', size=12, bold=True, color=GRAY7)
add_text(s, 0.80, 4.40, 9.00, 0.35,
         '프로젝트 기간  2026.05.15 — 2026.06.08  ·  v0.1.0 Build 1~3 + Polish',
         size=11, color=GRAY5)
add_text(s, 0.80, 4.78, 9.00, 0.35,
         '버전 2  ·  Build 3 이후 모바일 안정화/UX 개선 포함',
         size=10, italic=True, color=GRAY9)

# ---------------------------------------------------------------
# Slide 2 — Agenda
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, 'AGENDA', '오늘 발표 순서')
items = [
    ('01', '게임 소개',        '장르, 컨셉, 플랫폼, 핵심 루프'),
    ('02', '구현 범위',        '시스템과 콘텐츠 한눈에'),
    ('03', '개발 일정',        'Build 1 → 2 → 3 → Polish'),
    ('04', 'Deep Dive ①',     '웨이포인트 + 다중 경로 웨이브'),
    ('05', 'Deep Dive ②',     '타워 분기 진화 시스템'),
    ('06', 'Deep Dive ③',     'Build 3 이후 폴리시 & 버그 수정'),
    ('07', '느낀점 / Q & A',  '한 달의 회고와 시연'),
]
# 2 columns x 4 rows-ish — 7 items, do 4+3 layout
col_w = 4.40
col_x = [0.50, 5.10]
row_h = 0.85
for i, (num, title, desc) in enumerate(items):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = col_x[col]; y = 1.45 + row * row_h
    add_rect(s, x, y, col_w, 0.75, WHITE)
    add_rect(s, x, y, 0.08, 0.75, NAVY)
    add_text(s, x + 0.20, y + 0.05, 0.55, 0.40, num, size=20, bold=True, color=GOLD)
    add_text(s, x + 0.85, y + 0.06, 2.50, 0.30, title, size=13, bold=True, color=NAVY)
    add_text(s, x + 0.85, y + 0.38, 3.40, 0.30, desc, size=10, color=GRAY5)
add_footer(s, 1, TOTAL)

# ---------------------------------------------------------------
# Slide 3 — Game Overview
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, '01 · GAME OVERVIEW', '어떤 게임인가요?')
# Left content
add_text(s, 0.50, 1.50, 5.50, 0.55, '정통 타워 디펜스', size=20, bold=True, color=NAVY)
add_text(s, 0.50, 2.10, 5.50, 1.20,
         '지정 경로로 진격하는 적군이 목표 지점에 도달하지 못하도록,\n'
         '타워를 짓고 강화하며 막아내는 Unity 2D 모바일 게임.',
         size=12, color=GRAY7)
add_text(s, 0.50, 3.50, 5.50, 0.35, '핵심 재미', size=14, bold=True, color=GOLD)
add_text(s, 0.50, 3.85, 5.50, 1.40,
         '· 제한된 골드로 무엇을, 어디에, 어떤 순서로 강화할 것인가\n'
         '· 여러 경로에서 동시에 들어오는 압박을 어떻게 분산시킬 것인가\n'
         '· 위기 순간 — 특수능력을 언제 쓸 것인가',
         size=11, color=GRAY7)
# Right info card
add_rect(s, 6.40, 1.50, 3.10, 3.50, WHITE)
add_rect(s, 6.40, 1.50, 0.08, 3.50, NAVY)
add_text(s, 6.70, 1.62, 2.70, 0.30, 'PROJECT INFO', size=10, bold=True, color=GOLD)
info_rows = [
    ('장르',    '정통 Tower Defense'),
    ('엔진',    'Unity (2D)'),
    ('플랫폼',  'Mobile (Android APK)'),
    ('빌드',    'v0.1.0 Build 1~3+'),
    ('커밋',    '130+ commits'),
    ('기간',    '약 4주 + 폴리시'),
]
for i, (k, v) in enumerate(info_rows):
    y = 2.05 + i * 0.45
    add_text(s, 6.70, y, 0.90, 0.40, k, size=10, bold=True, color=NAVY)
    add_text(s, 7.55, y, 1.90, 0.40, v, size=10, color=GRAY7)
add_footer(s, 2, TOTAL)

# ---------------------------------------------------------------
# Slide 4 — Core Loop
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, '01 · GAME OVERVIEW', '한 판의 흐름 — Core Loop')
steps = [
    ('1', '웨이브 시작', '여러 경로에서 적이 진격'),
    ('2', '방어 구축',  '골드로 타워 건설 · 업그레이드'),
    ('3', '전투 / 위기', '특수능력으로 위기 돌파'),
    ('4', '보상 획득',  '골드 획득 · 다음 웨이브 준비'),
    ('5', '결과 판정',  '라이프 잔량으로 ★ 1~3 결정'),
]
step_w = 1.70
gap = 0.10
start_x = 0.50
for i, (n, t, d) in enumerate(steps):
    x = start_x + i * (step_w + gap)
    add_rect(s, x, 1.55, step_w, 2.30, WHITE)
    # number bubble
    add_rect(s, x + 0.55, 1.70, 0.60, 0.60, NAVY)
    add_text(s, x + 0.55, 1.73, 0.60, 0.55, n, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.05, 2.50, step_w - 0.10, 0.40, t, size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.05, 2.90, step_w - 0.10, 0.90, d, size=9, color=GRAY7,
             align=PP_ALIGN.CENTER)
# end conditions
add_rect(s, 0.50, 4.20, 9.00, 0.85, LBLUE)
add_text(s, 0.75, 4.30, 2.00, 0.30, '종료 조건', size=11, bold=True, color=NAVY)
add_text(s, 2.80, 4.30, 3.00, 0.30, 'Game Over  ·  라이프 ≤ 0',
         size=11, bold=True, color=RED)
add_text(s, 5.80, 4.30, 4.00, 0.30, 'Victory  ·  마지막 웨이브까지 전부 방어',
         size=11, bold=True, color=GREEN)
add_text(s, 0.75, 4.65, 9.00, 0.30,
         '결과 화면에서 라이프 비율로 ★1~3 평가 — 같은 클리어도 운영 효율로 차등화',
         size=9, italic=True, color=GRAY5)
add_footer(s, 3, TOTAL)

# ---------------------------------------------------------------
# Slide 5 — Systems
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, '02 · SCOPE', '구현한 시스템 한눈에')
sys_items = [
    ('타일맵 · 경로',     'Waypoint 기반 EnemyPath · 다중 경로'),
    ('웨이브 디렉터',     '병렬 SpawnEntry · 난이도 곡선 · 호출 보상'),
    ('타워 · 빌드 시스템', 'BuildSpot 라디얼 메뉴 · 분기 진화'),
    ('전투 · 페어 lock',  '보병/적 1:1 페어, 측면 교전, windup'),
    ('영웅 · 보병 · 적군', '근접/원거리/힐러/탱커/공중/보스'),
    ('특수능력',         '지원군 소환 · 용암 장판 (영역 지정)'),
    ('UI · HUD',         '골드/라이프/웨이브 · 일시정지 · ★ 결과'),
    ('모바일 안정화',     'Android 60fps · 화살 step-aware 명중 (NEW)'),
]
card_w = 2.22; card_h = 1.70
for i, (title, desc) in enumerate(sys_items):
    col = i % 4; row = i // 4
    x = 0.50 + col * (card_w + 0.08); y = 1.45 + row * (card_h + 0.12)
    add_rect(s, x, y, card_w, card_h, WHITE)
    is_new = '(NEW)' in desc
    bar_color = GOLD if is_new else NAVY
    add_rect(s, x, y, card_w, 0.08, bar_color)
    add_text(s, x + 0.18, y + 0.28, card_w - 0.36, 0.55, title,
             size=13, bold=True, color=NAVY)
    add_text(s, x + 0.18, y + 0.85, card_w - 0.36, 0.75,
             desc.replace(' (NEW)', ''),
             size=10, color=GRAY7)
    if is_new:
        add_text(s, x + card_w - 0.65, y + 0.18, 0.50, 0.25, 'NEW',
                 size=8, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
add_footer(s, 4, TOTAL)

# ---------------------------------------------------------------
# Slide 6 — Contents
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, '02 · SCOPE', '콘텐츠로 보는 구현 범위')
# Big stat row
stats = [
    ('3 + 4',  '타워 종류 + 분기 진화'),
    ('6+',    '적 유닛 (보스 포함)'),
    ('2 / 1', '특수능력 / 영웅'),
    ('130+',  '총 Git 커밋'),
]
for i, (num, label) in enumerate(stats):
    x = 0.50 + i * 2.30
    add_text(s, x, 1.45, 2.10, 0.75, num, size=36, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_text(s, x, 2.25, 2.10, 0.35, label, size=10, color=GRAY5,
             align=PP_ALIGN.CENTER)
# Content cards
add_text(s, 0.50, 2.80, 3.00, 0.30, 'CONTENTS', size=10, bold=True, color=GOLD)
contents = [
    ('타워',       'Archer · Mage · Barracks'),
    ('분기 진화',   'Ranger / Sniper / Pyromancer / Frost Mage'),
    ('지상 적',     '근접 · 원거리 · 힐러(Monk) · 탱커(Lancer)'),
    ('공중 · 보스', 'Bat · Boss (페이즈 시스템)'),
    ('아군',       '보병 · 영웅(Hero Knight) · 지원군'),
    ('특수능력',    'Reinforcement · LavaZone'),
]
for i, (k, v) in enumerate(contents):
    col = i % 2; row = i // 2
    x = 0.50 + col * 4.60; y = 3.20 + row * 0.55
    add_rect(s, x, y, 4.40, 0.45, WHITE)
    add_rect(s, x, y, 0.05, 0.45, NAVY)
    add_text(s, x + 0.18, y + 0.08, 1.20, 0.30, k, size=11, bold=True, color=NAVY)
    add_text(s, x + 1.50, y + 0.08, 2.80, 0.30, v, size=10, color=GRAY7)
add_footer(s, 5, TOTAL)

# ---------------------------------------------------------------
# Slide 7 — Timeline overview
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, '03 · TIMELINE', '약 4주 + 폴리시 — 4번의 사이클')
phases = [
    ('5/15-21', 'Build 1', '기반 시스템',
     ['프로젝트 셋업', '타일맵 · 경로', '타워 3종 · 웨이브', 'HUD · 골드 경제']),
    ('5/26-29', 'Build 2', '게임 루프',
     ['다중 경로 웨이브', '특수능력 2종', '신규 적 4종', '클리어/오버 UI']),
    ('6/02-08', 'Build 3', '콘텐츠 확장',
     ['영웅 유닛', '보스 페이즈', '메인메뉴 씬', '타워 분기 진화', '결과 화면 ★']),
    ('6/08~',  'Polish', 'QA & 안정화',
     ['모바일 화살 잔존 fix', '분기 진화 아이콘', '영웅 사망 모션 fix', 'Frost Mage 슬로우 틴트']),
]
col_w = 2.22; col_x_base = 0.50
for i, (date, code, focus, items) in enumerate(phases):
    x = col_x_base + i * (col_w + 0.08)
    is_polish = code == 'Polish'
    head_color = GOLD if is_polish else NAVY
    add_rect(s, x, 1.45, col_w, 3.55, WHITE)
    add_rect(s, x, 1.45, col_w, 0.50, head_color)
    add_text(s, x + 0.15, 1.50, col_w - 0.30, 0.20, date, size=8, bold=True,
             color=WHITE)
    add_text(s, x + 0.15, 1.70, col_w - 0.30, 0.25, code, size=14, bold=True,
             color=WHITE)
    add_text(s, x + 0.15, 2.05, col_w - 0.30, 0.30, focus, size=11, bold=True,
             color=NAVY)
    body = '\n'.join(f'· {it}' for it in items)
    add_text(s, x + 0.15, 2.45, col_w - 0.30, 2.50, body, size=9.5, color=GRAY7)
# bottom note
add_text(s, 0.50, 5.05, 9.00, 0.20,
         '각 빌드 종료마다 빌드 노트 + 결과 보고서 + 다음 빌드 계획 문서화. Polish 는 모바일 QA 발견 이슈 대응.',
         size=9, italic=True, color=GRAY5)
add_footer(s, 6, TOTAL)

# ---------------------------------------------------------------
# Slide 8 — Build highlights (4 columns)
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, '03 · TIMELINE', '빌드별 핵심 성과 — Build 1 → 2 → 3 → Polish')
cols = [
    ('BUILD 1', '기반 시스템', '최소 동작 게임 완성', NAVY,
     ['타워 3종 + Lv1→3 업글', '웨이포인트 경로', '단일 웨이브', '골드/라이프 HUD', '배럭 보병 자동 배치']),
    ('BUILD 2', '게임 루프 완성', '디펜스다운 깊이', GOLD,
     ['다중 경로 + 다중 웨이브', '특수능력 (지원군/용암)', '신규 적: 힐러 · 탱커', '랠리 포인트 지정']),
    ('BUILD 3', '콘텐츠 확장', '전략 다양성', NAVY,
     ['공중 적(Bat) + 보스 페이즈', '영웅 유닛 (배치형)', '분기 진화 (Ranger/Sniper)', '메인메뉴 + ★ 결과']),
    ('POLISH',  '모바일 안정화', '실 디바이스 QA 반영', GOLD,
     ['화살 step-aware 명중', '60fps 고정 (Android)', '분기 진화 아이콘 분리', '영웅 사망/부활 모션 fix', 'Frost Mage 슬로우 틴트']),
]
cw = 2.22
for i, (code, focus, desc, head_color, items) in enumerate(cols):
    x = 0.50 + i * (cw + 0.08)
    add_rect(s, x, 1.45, cw, 3.55, WHITE)
    add_rect(s, x, 1.45, cw, 0.65, head_color)
    add_text(s, x + 0.15, 1.50, cw - 0.30, 0.25, code, size=12, bold=True, color=WHITE)
    add_text(s, x + 0.15, 1.78, cw - 0.30, 0.25, focus, size=10.5, color=WHITE)
    add_text(s, x + 0.15, 2.20, cw - 0.30, 0.25, 'FOCUS', size=8, bold=True, color=GOLD)
    add_text(s, x + 0.15, 2.40, cw - 0.30, 0.30, desc, size=10, bold=True, color=NAVY)
    body = '\n'.join(f'· {it}' for it in items)
    add_text(s, x + 0.15, 2.80, cw - 0.30, 2.20, body, size=9, color=GRAY7)
add_footer(s, 7, TOTAL)

# ---------------------------------------------------------------
# Slide 9 — Deep Dive 1 intro
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 10, 5.625, NAVY)
add_rect(s, 0.80, 0.50, 0.18, 4.625, GOLD)
add_text(s, 1.20, 1.30, 8.00, 0.40, 'DEEP DIVE  ①', size=14, bold=True, color=GOLD)
add_text(s, 1.20, 1.80, 8.00, 1.30,
         '웨이포인트 경로 +\n다중 경로 웨이브 시스템',
         size=36, bold=True, color=WHITE)
add_text(s, 1.20, 3.60, 8.00, 0.90,
         '한 경로만 막던 디펜스에서, 동시에 들어오는\n여러 갈래를 어떻게 막을 것인가',
         size=14, italic=True, color=ICE)
add_text(s, 0.50, 5.26, 6.00, 0.25, 'Unity 2D Tower Defense  ·  김성수',
         size=9, color=ICE)
add_text(s, 8.80, 5.26, 0.70, 0.25, f'8 / {TOTAL}',
         size=9, color=ICE, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------
# Slide 10 — EnemyPath
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, 'DEEP DIVE ① · 경로', 'EnemyPath — 자식 Transform = 웨이포인트')
# Left: design ideas
add_text(s, 0.50, 1.50, 5.50, 0.35, '설계 아이디어', size=12, bold=True, color=GOLD)
ideas = [
    '· 자식 GameObject 순서 = 경로 순서',
    '· 에디터에서 자식 Transform 만 드래그하면 그대로 경로',
    '· Spawn = Waypoints[0], End = Waypoints[Count-1]',
    '· OnDrawGizmos 로 씬뷰에 점 + 선 표시 (디자이너 친화)',
    '· EnemyPath 여러 개 = 즉시 다중 경로 지원',
]
add_text(s, 0.50, 1.90, 5.50, 2.50, '\n'.join(ideas), size=11, color=GRAY7)
# Right: path diagram
add_rect(s, 6.30, 1.45, 3.20, 3.55, LBLUE)
add_text(s, 6.50, 1.55, 2.80, 0.30, '경로 구조 예시', size=11, bold=True, color=NAVY)
# Three paths visualization
path_y = [2.25, 3.00, 3.75]
path_label = ['L', 'M', 'R']
for i, py in enumerate(path_y):
    # spawn dot
    add_rect(s, 6.50, py, 0.18, 0.18, NAVY)
    # waypoints
    for j in range(3):
        wx = 6.85 + j * 0.55
        add_rect(s, wx, py + 0.05, 0.08, 0.08, GOLD)
    # end dot
    add_rect(s, 8.55, py, 0.18, 0.18, RED)
    add_text(s, 8.85, py - 0.03, 0.40, 0.30, path_label[i], size=11, bold=True, color=NAVY)
add_text(s, 6.50, 4.55, 2.80, 0.30,
         '●스폰  ◆웨이포인트  ●골인', size=9, color=GRAY5)
add_footer(s, 9, TOTAL)

# ---------------------------------------------------------------
# Slide 11 — WaveData
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, 'DEEP DIVE ① · 웨이브', 'WaveData — startOffset 기반 병렬 스폰')
# Top description
add_rect(s, 0.50, 1.50, 9.00, 1.10, WHITE)
add_rect(s, 0.50, 1.50, 0.08, 1.10, NAVY)
add_text(s, 0.70, 1.62, 8.70, 0.95,
         '한 웨이브 안의 SpawnEntry 들은 startOffset 시점에 "동시에" 시작된다.\n'
         '같은 시간 = 좌우 동시 발진, 시간차 = 시간차 진입.\n'
         '다중 경로 압박을 데이터만으로 표현.',
         size=11, color=GRAY7)
# Director details
add_rect(s, 0.50, 2.80, 9.00, 2.20, LBLUE)
add_text(s, 0.75, 2.95, 8.50, 0.30, 'WaveDirector', size=14, bold=True, color=NAVY)
director_items = [
    ('• pathId → EnemySpawner 라우팅',     '경로별 SpawnEntry 분리, 동일 웨이브에서 다중 경로 처리'),
    ('• AnimationCurve 난이도 곡선',       'count / interval / hp 곡선으로 후반 가중 증폭'),
    ('• 다음 웨이브 일찍 호출 → 능력 쿨감', '절약된 시간을 특수능력 쿨다운 보상으로 환원'),
]
for i, (k, v) in enumerate(director_items):
    y = 3.35 + i * 0.55
    add_text(s, 0.75, y, 4.20, 0.30, k, size=11, bold=True, color=NAVY)
    add_text(s, 5.00, y, 4.50, 0.30, v, size=10, color=GRAY7)
add_footer(s, 10, TOTAL)

# ---------------------------------------------------------------
# Slide 12 — Deep Dive 2 intro
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 10, 5.625, NAVY)
add_rect(s, 0.80, 0.50, 0.18, 4.625, GOLD)
add_text(s, 1.20, 1.30, 8.00, 0.40, 'DEEP DIVE  ②', size=14, bold=True, color=GOLD)
add_text(s, 1.20, 1.80, 8.00, 1.30,
         '타워 분기 진화 시스템', size=36, bold=True, color=WHITE)
add_text(s, 1.20, 3.30, 8.00, 1.30,
         'Lv1 → Lv2 → Lv3 까지 단일 진화.\n'
         'Lv3 에서 두 갈래로 분기 — 같은 타워도 다른 전략.',
         size=14, italic=True, color=ICE)
add_text(s, 0.50, 5.26, 6.00, 0.25, 'Unity 2D Tower Defense  ·  김성수',
         size=9, color=ICE)
add_text(s, 8.80, 5.26, 0.70, 0.25, f'11 / {TOTAL}',
         size=9, color=ICE, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------
# Slide 13 — BuildingData architecture
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, 'DEEP DIVE ② · 아키텍처',
                   'BuildingData.nextBranches — 데이터 한 줄로 분기')
# Left card
add_rect(s, 0.50, 1.45, 4.70, 3.55, WHITE)
add_rect(s, 0.50, 1.45, 0.08, 3.55, NAVY)
add_text(s, 0.75, 1.55, 4.40, 0.30, '핵심 아이디어', size=11, bold=True, color=GOLD)
add_text(s, 0.75, 1.85, 4.40, 0.35, 'ScriptableObject 한 필드의 추가',
         size=13, bold=True, color=NAVY)
# Code block
add_rect(s, 0.75, 2.30, 4.30, 1.40, CODE)
code = ('public BuildingData nextUpgrade;\n'
        'public BuildingData[] nextBranches;\n'
        'public Sprite upgradeSlotIcon;  // NEW\n'
        '// 분기 1+ 이면 nextBranches 우선\n'
        '// 슬롯 아이콘은 분기별 시각 구분용')
add_text(s, 0.90, 2.38, 4.00, 1.28, code, size=9.5, color=ICE, font='Consolas')
add_text(s, 0.75, 3.85, 4.40, 1.15,
         '· 기존 단일 업그레이드 자산은 그대로 동작 (회귀 없음)\n'
         '· 분기를 쓸 타워만 nextBranches 채우면 끝\n'
         '· 새 upgradeSlotIcon 으로 분기 슬롯 시각화 (Polish)',
         size=10, color=GRAY7)
# Right card - radial menu
add_rect(s, 5.40, 1.45, 4.10, 3.55, LBLUE)
add_text(s, 5.60, 1.55, 3.80, 0.30, '라디얼 메뉴 분기 표시',
         size=11, bold=True, color=NAVY)
# Center tower
cx, cy = 7.30, 3.00
add_rect(s, cx - 0.55, cy - 0.55, 1.10, 1.10, NAVY)
add_text(s, cx - 0.55, cy - 0.55, 1.10, 1.10, 'TOWER',
         size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Branch A (upper left)
add_rect(s, cx - 1.40, cy - 1.40, 0.84, 0.84, GOLD)
add_text(s, cx - 1.40, cy - 1.40, 0.84, 0.84, '분기 A',
         size=10, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Branch B (upper right)
add_rect(s, cx + 0.55, cy - 1.40, 0.84, 0.84, GOLD)
add_text(s, cx + 0.55, cy - 1.40, 0.84, 0.84, '분기 B',
         size=10, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Sell (bottom)
add_rect(s, cx - 0.30, cy + 0.90, 0.60, 0.60, RED)
add_text(s, cx - 0.30, cy + 0.90, 0.60, 0.60, 'SELL',
         size=9, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 5.60, 4.70, 3.80, 0.30,
         '후보 1개 → 12시 · 2개 → 11/1시 · 3개+ 균등',
         size=8.5, color=GRAY5, align=PP_ALIGN.CENTER)
add_footer(s, 12, TOTAL)

# ---------------------------------------------------------------
# Slide 14 — Branch details
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, 'DEEP DIVE ② · 분기',
                   '두 갈래, 네 가지 전략 — Archer & Mage')
# Archer card
add_rect(s, 0.50, 1.45, 4.40, 3.30, WHITE)
add_rect(s, 0.50, 1.45, 4.40, 0.55, NAVY)
add_text(s, 0.70, 1.52, 4.00, 0.40, 'ARCHER  ·  Lv3', size=13, bold=True, color=WHITE)
# Ranger
add_rect(s, 0.65, 2.15, 4.10, 1.20, LBLUE)
add_text(s, 0.85, 2.22, 3.80, 0.30, 'Ranger', size=13, bold=True, color=NAVY)
add_text(s, 0.85, 2.50, 3.80, 0.30, '빠른 다단 사격', size=10.5, bold=True, color=GOLD)
add_text(s, 0.85, 2.78, 3.80, 0.50,
         '공격속도↑ — 다수 처리에 유리한 DPS 분포형',
         size=10, color=GRAY7)
# Sniper
add_rect(s, 0.65, 3.45, 4.10, 1.20, LBLUE)
add_text(s, 0.85, 3.52, 3.80, 0.30, 'Sniper', size=13, bold=True, color=NAVY)
add_text(s, 0.85, 3.80, 3.80, 0.30, '고데미지 저격', size=10.5, bold=True, color=GOLD)
add_text(s, 0.85, 4.08, 3.80, 0.50,
         '단일 고데미지 — 탱커/보스 단숨에 깎는 단발형',
         size=10, color=GRAY7)
# Mage card
add_rect(s, 5.10, 1.45, 4.40, 3.30, WHITE)
add_rect(s, 5.10, 1.45, 4.40, 0.55, NAVY)
add_text(s, 5.30, 1.52, 4.00, 0.40, 'MAGE  ·  Lv3', size=13, bold=True, color=WHITE)
# Pyromancer
add_rect(s, 5.25, 2.15, 4.10, 1.20, LBLUE)
add_text(s, 5.45, 2.22, 3.80, 0.30, 'Pyromancer', size=13, bold=True, color=NAVY)
add_text(s, 5.45, 2.50, 3.80, 0.30, '광역 (AoE)', size=10.5, bold=True, color=GOLD)
add_text(s, 5.45, 2.78, 3.80, 0.50,
         'splashRadius — 명중 지점 반경 내 적 전원 피해',
         size=10, color=GRAY7)
# Frost Mage
add_rect(s, 5.25, 3.45, 4.10, 1.20, ICE)
add_text(s, 5.45, 3.52, 3.80, 0.30, 'Frost Mage', size=13, bold=True, color=NAVY)
add_text(s, 5.45, 3.80, 3.80, 0.30, '둔화 (Slow) + 청색 틴트',
         size=10.5, bold=True, color=NAVY)
add_text(s, 5.45, 4.08, 3.80, 0.50,
         'slowAmount + slowDuration · 피격 적이 시각적으로도 얼어붙음',
         size=10, color=GRAY7)
add_text(s, 0.50, 4.85, 9.00, 0.30,
         '같은 자원 · 같은 위치 · 다른 전략 — 분기 선택만으로 운영이 달라진다',
         size=10, italic=True, color=GRAY5, align=PP_ALIGN.CENTER)
add_footer(s, 13, TOTAL)

# ---------------------------------------------------------------
# Slide 15 — Deep Dive 3 intro (Polish)
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 10, 5.625, NAVY)
add_rect(s, 0.80, 0.50, 0.18, 4.625, GOLD)
add_text(s, 1.20, 1.30, 8.00, 0.40, 'DEEP DIVE  ③', size=14, bold=True, color=GOLD)
add_text(s, 1.20, 1.80, 8.00, 1.30,
         'Build 3+ Polish\n— 모바일 QA & 사용성 개선',
         size=32, bold=True, color=WHITE)
add_text(s, 1.20, 3.50, 8.00, 1.30,
         '실제 Android 디바이스에서 발견된 4가지 이슈와\n'
         '코드/데이터로 해결한 과정',
         size=14, italic=True, color=ICE)
add_text(s, 0.50, 5.26, 6.00, 0.25, 'Unity 2D Tower Defense  ·  김성수',
         size=9, color=ICE)
add_text(s, 8.80, 5.26, 0.70, 0.25, f'14 / {TOTAL}',
         size=9, color=ICE, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------
# Slide 16 — Polish details (4 fixes)
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, 'DEEP DIVE ③ · 폴리시',
                   '4가지 핵심 수정 — 데이터로 본 차이')
fixes = [
    ('01', '모바일 화살 잔존',
     'Android 기본 30fps 캡 → 화살이 적을 건너뜀',
     '· Application.targetFrameRate = 60 부트스트랩\n'
     '· step >= 남은거리 명중 판정 (Arrow / Magic)',
     NAVY),
    ('02', '분기 진화 아이콘 분리',
     '단일 업글과 분기 슬롯이 시각적으로 동일',
     '· BuildingData.upgradeSlotIcon 필드 추가\n'
     '· 공용 evolveIcon 폴백 (분기/단일 자동 구분)',
     GOLD),
    ('03', '영웅 사망/부활 모션 꼬임',
     'Attack 중 사망 시 부활하면 자세 잔재',
     '· Animator.Play("Die"/"Idle", 0, 0f) 강제 점프\n'
     '· 트랜지션 그래프 결함(Attack→Die 부재) 우회',
     NAVY),
    ('04', 'Frost Mage 슬로우 시각',
     '슬로우 걸린 적인지 한눈에 안 보임',
     '· 적 SpriteRenderer 들을 청색 톤으로 일괄 틴트\n'
     '· 만료 시 원본 복원 (페이즈 깜빡임과 공존)',
     GOLD),
]
cw = 4.40; ch = 1.65
for i, (num, title, problem, fix, accent) in enumerate(fixes):
    col = i % 2; row = i // 2
    x = 0.50 + col * (cw + 0.20); y = 1.45 + row * (ch + 0.15)
    add_rect(s, x, y, cw, ch, WHITE)
    add_rect(s, x, y, 0.08, ch, accent)
    add_text(s, x + 0.22, y + 0.08, 0.50, 0.30, num, size=16, bold=True, color=accent)
    add_text(s, x + 0.80, y + 0.10, 3.50, 0.30, title, size=12, bold=True, color=NAVY)
    add_text(s, x + 0.22, y + 0.48, cw - 0.40, 0.30,
             f'문제 · {problem}', size=9, italic=True, color=RED)
    add_text(s, x + 0.22, y + 0.78, cw - 0.40, 0.80,
             fix, size=9.5, color=GRAY7)
add_footer(s, 15, TOTAL)

# ---------------------------------------------------------------
# Slide 17 — Retrospective
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_section_header(s, '04 · RETROSPECTIVE', '한 달, 그리고 남은 것들')
retro = [
    ('GOOD', '데이터 주도 설계의 힘',
     'ScriptableObject(EnemyData / TowerData / WaveData / BuildingData)로 콘텐츠 자산화. '
     '새 적/타워/분기를 코드 없이 추가.', GREEN),
    ('GOOD', '빌드 사이클을 반복한 효과',
     'Build 1~3 + Polish 마다 빌드 노트 · 결과 보고서 · 다음 빌드 계획 문서화. '
     '우선순위가 흐려지지 않음.', GREEN),
    ('HARD', '전투 디테일 + 모바일 QA',
     '보병/적 1:1 페어 lock, 측면 교전, windup, 그리고 30fps 환경에서만 보이는 '
     '오버슈트 — 가장 많은 패치를 부른 영역.', GOLD),
    ('NEXT', '다음에 해보고 싶은 것',
     'FindObjectsByType 매 프레임 호출 → EnemyManager 등록제로 교체, '
     '사운드/이펙트, 스테이지 다양화.', NAVY),
]
cw = 4.40; ch = 1.70
for i, (tag, title, desc, color) in enumerate(retro):
    col = i % 2; row = i // 2
    x = 0.50 + col * (cw + 0.20); y = 1.45 + row * (ch + 0.10)
    add_rect(s, x, y, cw, ch, WHITE)
    add_rect(s, x, y, 0.08, ch, color)
    add_text(s, x + 0.22, y + 0.10, 1.00, 0.30, tag, size=10, bold=True, color=color)
    add_text(s, x + 0.22, y + 0.40, cw - 0.40, 0.40, title,
             size=13, bold=True, color=NAVY)
    add_text(s, x + 0.22, y + 0.85, cw - 0.40, 0.80, desc, size=10, color=GRAY7)
add_footer(s, 16, TOTAL)

# ---------------------------------------------------------------
# Slide 18 — Demo Video (dedicated)
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
# 흰 배경. 중앙에 큰 비디오 플레이스홀더 프레임.
add_rect(s, 0, 0, 10, 5.625, WHITE)
# 섹션 헤더 (다른 슬라이드와 동일 톤)
add_text(s, 0.50, 0.32, 9.00, 0.30, '05 · DEMO', size=11, bold=True, color=GOLD)
add_text(s, 0.50, 0.58, 9.00, 0.70, '시연 영상', size=24, bold=True, color=NAVY)

# 비디오 플레이스홀더 — 검정 화면 + 가운데 ▶ 버튼 + 16:9 비율 유지.
# 슬라이드 영역 안에서 최대한 크게 (좌우 0.5 마진, 위아래 푸터 공간 제외).
# 가용 폭 9.0", 가용 높이 약 3.9" → 16:9 맞추면 6.93" x 3.9" 또는 6.4" x 3.6". 6.93×3.9 채택.
vw, vh = 6.93, 3.9
vx = (10 - vw) / 2
vy = 1.45
add_rect(s, vx, vy, vw, vh, RGBColor(0x12, 0x12, 0x12))
# 얇은 골드 테두리 느낌 — 좌상단 모서리 강조 (decorative line 은 금지지만 작은 brand marker 는 OK)
add_rect(s, vx, vy, 0.18, 0.05, GOLD)
add_rect(s, vx, vy, 0.05, 0.18, GOLD)
add_rect(s, vx + vw - 0.18, vy + vh - 0.05, 0.18, 0.05, GOLD)
add_rect(s, vx + vw - 0.05, vy + vh - 0.18, 0.05, 0.18, GOLD)
# 가운데 ▶ 버튼 — 원 + 삼각형 느낌. ISOSCELES_TRIANGLE 도형 사용.
cx_btn = vx + vw / 2
cy_btn = vy + vh / 2
btn_r = 0.55
btn = s.shapes.add_shape(MSO_SHAPE.OVAL,
                         Inches(cx_btn - btn_r), Inches(cy_btn - btn_r),
                         Inches(btn_r * 2), Inches(btn_r * 2))
btn.fill.solid(); btn.fill.fore_color.rgb = GOLD
btn.line.fill.background()
btn.shadow.inherit = False
# 재생 삼각형 — 오른쪽 가리키도록 회전 90도.
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                         Inches(cx_btn - 0.20), Inches(cy_btn - 0.25),
                         Inches(0.45), Inches(0.50))
# RIGHT_TRIANGLE 은 직각이라 어울리지 않으므로 isosceles 로 교체 후 회전.
sp_xml = tri._element
sp_xml.getparent().remove(sp_xml)
play = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                          Inches(cx_btn - 0.22), Inches(cy_btn - 0.25),
                          Inches(0.44), Inches(0.50))
play.rotation = 90
play.fill.solid(); play.fill.fore_color.rgb = WHITE
play.line.fill.background()
play.shadow.inherit = False

# 비디오 라벨 — 우상단 작게.
add_text(s, vx + 0.20, vy + 0.20, 2.50, 0.30,
         'GAMEPLAY DEMO', size=10, bold=True, color=GOLD)
# 비디오 메타 — 좌하단.
add_text(s, vx + 0.20, vy + vh - 0.45, 6.50, 0.30,
         'v0.1.0  ·  Android APK  ·  Build 3 + Polish',
         size=9, italic=True, color=ICE)

# 슬라이드 하단 안내 — 푸터 자리에 발표자 노트 형식으로.
add_text(s, 0.50, 5.40, 9.00, 0.20,
         '발표 시 이 페이지에서 영상 재생 — 영상 파일을 PPT 에 임베드하거나 외부 플레이어로 전환',
         size=9, italic=True, color=GRAY9, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------
# Slide 19 — Thanks / Q & A
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 10, 5.625, NAVY)
add_text(s, 0.80, 0.50, 9.00, 0.40, 'Q & A', size=12, bold=True, color=GOLD)
add_text(s, 0.80, 1.50, 9.00, 1.10, 'THANK YOU',
         size=54, bold=True, color=WHITE)
add_rect(s, 0.80, 2.70, 1.60, 0.05, GOLD)
add_text(s, 0.80, 2.95, 9.00, 0.45,
         '발표를 들어주셔서 감사합니다.', size=18, color=WHITE)
add_text(s, 0.80, 3.45, 9.00, 0.45,
         '자유롭게 질문해 주세요.', size=14, color=ICE)
add_text(s, 0.80, 4.50, 9.00, 0.30,
         '김성수  ·  Unity 2D Tower Defense  ·  v0.1.0 Build 3 + Polish',
         size=11, italic=True, color=ICE)
add_text(s, 0.80, 4.85, 9.00, 0.30,
         'rlatjdtn620@gmail.com',
         size=10, color=GRAY9)

import os
out = 'SwordRush 발표자료 v2.pptx'
try:
    prs.save(out)
    print(f'Saved: {out}')
except PermissionError:
    fallback = 'SwordRush 발표자료 v2_new.pptx'
    prs.save(fallback)
    print(f'기존 {out} 이 열려있어 {fallback} 으로 저장. 기존 파일 닫은 뒤 덮어쓰기 또는 이 파일을 사용.')
print(f'Total slides: {len(prs.slides)}')
